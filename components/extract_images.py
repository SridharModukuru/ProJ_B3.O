import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from io import BytesIO

def extract_text_and_images_from_pptx(pptx_file, output_folder):
    prs = Presentation(pptx_file)
    extracted_data = []

    for i, slide in enumerate(prs.slides):
        slide_data = {'text': '', 'images': []}
        # Extract text from slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_data['text'] += shape.text + '\n'

            # Extract images from slide
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_stream = BytesIO(image_bytes)
                image_obj = Image.open(image_stream)

                # Save the image
                image_filename = f'slide_{i + 1}_image.png'
                image_path = os.path.join(output_folder, image_filename)
                image_obj.save(image_path)
                slide_data['images'].append(image_path)

        extracted_data.append(slide_data)

    return extracted_data

# extract_text_and_images_from_pptx(r'C:\Users\Sridhar\_EMPTY FOLDERS\ProJ_B3.O\uploaded_files\Global-Warming.pptx',r'C:\Users\Sridhar\_EMPTY FOLDERS\ProJ_B3.O\slide_images'))

# print()
# print('*************************')

# print(extract_text_and_images_from_pptx(r'C:\Users\Sridhar\_EMPTY FOLDERS\ProJ_B3.O\uploaded_files\Global-Warming.pptx',r'C:\Users\Sridhar\_EMPTY FOLDERS\ProJ_B3.O\slide_images')[0])